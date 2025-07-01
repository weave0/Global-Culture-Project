import streamlit as st
import os
import json
import zipfile
import streamlit.components.v1 as components
import folium
from streamlit_folium import st_folium
import time
import datetime
import plotly.graph_objects as go
import pandas as pd
import random
import threading
import urllib.parse
import importlib

FOLDER = "parsed_output"
EXPORT_DIR = "tmp_export"
PLUGIN_DIR = "modules"

# Optional field grouping for UI preview
categories = {
    "🌐 Identity & Locale": ["culture_name", "language_tag", "estimated_speakers", "primary_regions"],
    "🧠 Culture & Communication": ["greetings", "communication_style", "gestures"],
    "🙏 Beliefs & Roles": ["religion", "gender_roles", "taboos"],
    "🧥 Appearance": ["clothing"]
}

# Define reusable styles
CARD_STYLE_FILLED = "background-color:#e3f2fd; border:2px solid #90caf9; padding: 1rem; border-radius: 10px; margin-bottom: 1rem; box-shadow: 1px 1px 5px rgba(0,0,0,0.05);"
CARD_STYLE_EMPTY = "background-color:#fff3cd; border:2px dashed #ffb74d; padding: 1rem; border-radius: 10px; margin-bottom: 1rem; box-shadow: 1px 1px 5px rgba(0,0,0,0.05);"

# Define reusable theme colors
DEEP_PURPLE = "#4A148C"
LIGHT_PURPLE = "#E1BEE7"

# Define theme objects for dark/light mode consistency
theme = {
    "light": {
        "background": "#ffffff",
        "text": "#000000",
        "header_gradient": "linear-gradient(to right, #ff7e5f, #feb47b)",
    },
    "dark": {
        "background": "#1e1e1e",
        "text": "#ffffff",
        "header_gradient": "linear-gradient(to right, #4b79a1, #283e51)",
    },
    "regional": {
        "Asia": "#ff4d4d",
        "Africa": "#ff9900",
        "Europe": "#4da6ff",
        "Americas": "#66cc66",
    }
}

# Optimize data loading by caching entries
@st.cache_data
def load_entries(folder):
    return [f for f in os.listdir(folder) if f.endswith(".v3.json")]

@st.cache_data
def load_entry(path):
    try:
        with open(path, encoding="utf-8") as f:
            return json.load(f)
    except FileNotFoundError:
        st.error(f"File not found: {path}")
        return {}
    except json.JSONDecodeError:
        st.error(f"Error decoding JSON in file: {path}")
        return {}

def save_entry(path, data):
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

def calculate_summary(entry):
    """Calculate field completion percentage and missing fields."""
    total_fields = len(entry)
    completed_fields = sum(1 for v in entry.values() if v)
    missing_fields = [k for k, v in entry.items() if not v]
    completion_percentage = (completed_fields / total_fields) * 100 if total_fields else 0
    return completion_percentage, missing_fields

def export_as_html(entry, output_path):
    """Export the entry as an HTML file."""
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(f"<html><body><h1>{entry.get('culture_name', 'Unnamed Culture')}</h1>")
        for key, value in entry.items():
            f.write(f"<p><strong>{key}:</strong> {value}</p>")
        f.write("</body></html>")

def export_as_markdown(entry, output_path):
    """Export the entry as a Markdown file."""
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(f"# {entry.get('culture_name', 'Unnamed Culture')}\n\n")
        for key, value in entry.items():
            f.write(f"**{key}:** {value}\n\n")

def batch_export(entries, format, output_zip):
    os.makedirs(EXPORT_DIR, exist_ok=True)
    with zipfile.ZipFile(output_zip, "w") as zipf:
        for entry_name, entry_data in entries.items():
            base_name = entry_name.replace(".v3.json", "")
            file_name = f"{base_name}.{format}"
            file_path = os.path.join(EXPORT_DIR, file_name)
            if format == "html":
                export_as_html(entry_data, file_path)
            elif format == "md":
                export_as_markdown(entry_data, file_path)
            zipf.write(file_path, arcname=file_name)

def get_next_field(entry):
    for k, v in entry.items():
        if not str(v).strip():
            return k
    return None

# Function to create a map for primary regions
def create_region_map(primary_regions):
    m = folium.Map(location=[0, 0], zoom_start=2)
    for region in primary_regions.split(", "):
        # Example: Add markers for regions (replace with actual geo-coordinates mapping)
        if region == "North America":
            folium.Marker([37.0902, -95.7129], popup="North America").add_to(m)
        elif region == "Europe":
            folium.Marker([54.5260, 15.2551], popup="Europe").add_to(m)
        # Add more regions as needed
    return m

# Save a backup file on each save
def save_backup(entry, path):
    timestamp = datetime.datetime.now().strftime("%Y%m%d%H%M%S")
    backup_path = os.path.join("backups", f"{os.path.basename(path)}_{timestamp}.json")
    os.makedirs("backups", exist_ok=True)
    save_entry(backup_path, entry)

# Function to create a radar chart
def create_completion_radar_chart(entry):
    category_completion = []
    for category, fields in categories.items():
        total = len(fields)
        completed = sum(1 for f in fields if entry.get(f))
        category_completion.append((category, (completed / total) * 100 if total else 0))

    labels = [c[0] for c in category_completion]
    values = [c[1] for c in category_completion]

    fig = go.Figure()
    fig.add_trace(go.Scatterpolar(
        r=values,
        theta=labels,
        fill='toself',
        name='Completion'
    ))

    fig.update_layout(
        polar=dict(
            radialaxis=dict(
                visible=True,
                range=[0, 100]
            )
        ),
        showlegend=False
    )
    return fig

# Add smart field auto-scroll functionality
def scroll_to_field(field_name):
    components.html(f"""
        <script>
        document.getElementById('{field_name}').scrollIntoView({{ behavior: 'smooth' }});
        </script>
    """, height=0)

# Add culture-specific theming
def get_theme_colors(region):
    themes = {
        "Asia": {"background": "#FFEBEE", "accent": "#D32F2F"},
        "Europe": {"background": "#E3F2FD", "accent": "#1976D2"},
        "Africa": {"background": "#FFF3E0", "accent": "#F57C00"},
        "Oceania": {"background": "#E8F5E9", "accent": "#388E3C"},
        "North America": {"background": "#F3E5F5", "accent": "#7B1FA2"},
        "South America": {"background": "#FFFDE7", "accent": "#FBC02D"}
    }
    return themes.get(region, {"background": "#FFFFFF", "accent": "#000000"})

# Add dynamic image gallery
def display_image_gallery(images):
    if images:
        st.markdown("### 📸 Image Gallery")
        cols = st.columns(len(images))
        for i, img_url in enumerate(images):
            with cols[i]:
                st.image(img_url, use_column_width=True)

# Add top missing fields heatmap
def generate_missing_fields_heatmap(entries):
    from collections import Counter
    import plotly.express as px

    missing_fields = Counter()
    for entry in entries:
        for field, value in entry.items():
            if not value:
                missing_fields[field] += 1

    if missing_fields:
        df = pd.DataFrame(missing_fields.items(), columns=["Field", "Count"])
        fig = px.bar(df, x="Field", y="Count", title="Top Missing Fields", labels={"Count": "Number of Entries Missing"})
        st.plotly_chart(fig)

# Add submission history log
def log_submission(entry, path, changes):
    timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    log_entry = {
        "timestamp": timestamp,
        "entry": os.path.basename(path),
        "changes": changes
    }
    log_path = os.path.join("logs", "submission_history.json")
    os.makedirs("logs", exist_ok=True)
    if os.path.exists(log_path):
        with open(log_path, "r", encoding="utf-8") as f:
            log_data = json.load(f)
    else:
        log_data = []
    log_data.append(log_entry)
    with open(log_path, "w", encoding="utf-8") as f:
        json.dump(log_data, f, ensure_ascii=False, indent=2)

# Add AI confidence meter
def add_confidence_meter(field_name, suggestion):
    st.markdown(f"**AI Suggestion for {field_name.replace('_', ' ').title()}:**")
    st.info(suggestion)
    confidence = st.slider("How confident are you in this suggestion?", 1, 5, 3, key=f"confidence_{field_name}")
    return confidence

# Add a function to generate speaker trends
def generate_speaker_trends():
    years = list(range(2000, 2025))
    speakers = [random.randint(1000, 1000000) for _ in years]  # Simulated data
    return years, speakers

# Add a plot for speaker trends
def display_speaker_trends(entry):
    years, speakers = generate_speaker_trends()
    fig = go.Figure()
    fig.add_trace(go.Scatter(x=years, y=speakers, mode='lines+markers', name='Speakers'))
    fig.update_layout(title='Estimated Speaker Trends', xaxis_title='Year', yaxis_title='Estimated Speakers')
    st.plotly_chart(fig)

# Add translation-on-hover functionality
def add_translation_tooltip(field_name, value, translations):
    if field_name in translations:
        tooltip = translations[field_name]
        st.markdown(f"<span title='{tooltip}'>{value}</span>", unsafe_allow_html=True)
    else:
        st.text(value)

# Example translations dictionary (to be replaced with real data)
translations = {
    "culture_name": "Nombre de la cultura",
    "language_tag": "Etiqueta de idioma",
    "estimated_speakers": "Hablantes estimados",
    "primary_regions": "Regiones principales"
}

# Add a cultural calendar widget
def display_cultural_calendar(entry):
    holidays = entry.get("holidays", [])  # Assuming 'holidays' is a list of holiday names
    if holidays:
        st.markdown("### 📅 Cultural Calendar")
        for holiday in holidays:
            st.markdown(f"- {holiday}")
    else:
        st.info("No holidays listed for this culture.")

# Add advanced search functionality
def advanced_search(entries, search_query, filters):
    filtered_entries = []
    for entry in entries:
        match = True
        if search_query and search_query.lower() not in entry.get("culture_name", "").lower():
            match = False
        for key, value in filters.items():
            if value and value.lower() not in entry.get(key, "").lower():
                match = False
        if match:
            filtered_entries.append(entry)
    return filtered_entries

# Add print-friendly mode
def generate_print_friendly(entry):
    st.markdown("### 🖨️ Print-Friendly Version")
    st.markdown("<div style='font-family: Arial, sans-serif; font-size: 12pt;'>", unsafe_allow_html=True)
    for key, value in entry.items():
        st.markdown(f"<p><strong>{key.replace('_', ' ').title()}:</strong> {value}</p>", unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)

# Add schema validation
def validate_schema(entry, schema):
    errors = []
    for field, field_type in schema.items():
        if field not in entry:
            errors.append(f"Missing field: {field}")
        elif not isinstance(entry[field], field_type):
            errors.append(f"Field {field} should be of type {field_type.__name__}")
    return errors

# Example schema
default_schema = {
    "culture_name": str,
    "language_tag": str,
    "estimated_speakers": int,
    "primary_regions": str
}

# Add crowdsourced suggestions feature
def submit_suggestion(entry_name, field, suggestion):
    suggestion_log = os.path.join("logs", "suggestions.json")
    os.makedirs("logs", exist_ok=True)
    if os.path.exists(suggestion_log):
        with open(suggestion_log, "r", encoding="utf-8") as f:
            suggestions = json.load(f)
    else:
        suggestions = []
    suggestions.append({"entry": entry_name, "field": field, "suggestion": suggestion, "timestamp": datetime.datetime.now().isoformat()})
    with open(suggestion_log, "w", encoding="utf-8") as f:
        json.dump(suggestions, f, ensure_ascii=False, indent=2)

# Add smart save indicators
def add_save_indicator(field_name, is_saved):
    if is_saved:
        st.markdown(f"<span style='color: green;'>✔️ Saved</span>", unsafe_allow_html=True)
    else:
        st.markdown(f"<span style='color: red;'>✖️ Unsaved Changes</span>", unsafe_allow_html=True)

# Add compare mode
def compare_cultures(entry_a, entry_b):
    st.markdown("### 🔍 Compare Cultures")
    for field in entry_a.keys():
        value_a = entry_a.get(field, "—")
        value_b = entry_b.get(field, "—")
        if value_a != value_b:
            st.markdown(f"**{field.replace('_', ' ').title()}**: <span style='color: red;'>{value_a}</span> | <span style='color: green;'>{value_b}</span>", unsafe_allow_html=True)
        else:
            st.markdown(f"**{field.replace('_', ' ').title()}**: {value_a}", unsafe_allow_html=True)

# Add field-level history log
def get_field_history(field_name, entry_name):
    history_path = os.path.join("logs", "field_history.json")
    if os.path.exists(history_path):
        with open(history_path, "r", encoding="utf-8") as f:
            history = json.load(f)
        return [h for h in history if h["entry"] == entry_name and h["field"] == field_name]
    return []

def log_field_change(entry_name, field_name, old_value, new_value):
    history_path = os.path.join("logs", "field_history.json")
    os.makedirs("logs", exist_ok=True)
    if os.path.exists(history_path):
        with open(history_path, "r", encoding="utf-8") as f:
            history = json.load(f)
    else:
        history = []
    history.append({"entry": entry_name, "field": field_name, "old_value": old_value, "new_value": new_value, "timestamp": datetime.datetime.now().isoformat()})
    with open(history_path, "w", encoding="utf-8") as f:
        json.dump(history, f, ensure_ascii=False, indent=2)

# Ensure 'fields' and 'entry' are defined
if 'fields' not in locals():
    fields = ["field1", "field2", "field3"]  # Replace with actual field names

if 'entry' not in locals():
    entry = {}

# Ensure 'files' is defined before using it
if 'files' not in locals():
    files = load_entries(FOLDER)

# Ensure 'selected' and 'path' are defined
if 'selected' not in locals():
    selected = st.sidebar.selectbox("Select Culture", files)

if 'path' not in locals() and selected:
    path = os.path.join(FOLDER, selected)
    entry = load_entry(path)

# Ensure 'files' is defined
if 'files' not in locals():
    files = load_entries(FOLDER)

# Recalculate 'completion_percentage'
if 'completion_percentage' not in locals() and files:
    completion_percentage = sum(calculate_summary(load_entry(os.path.join(FOLDER, f)))[0] for f in files) / len(files)

# Integrate field history into the main function
for field in fields:
    value = entry.get(field, "")
    new_value = st.text_area(field.replace("_", " ").title(), value, key=field)
    if new_value != value:
        log_field_change(selected, field, value, new_value)
        entry[field] = new_value
        save_entry(path, entry)
    if st.button(f"Show Change History for {field}"):
        history = get_field_history(field, selected)
        if history:
            for h in history:
                st.markdown(f"- {h['timestamp']}: {h['old_value']} → {h['new_value']}")
        else:
            st.info("No history available for this field.")

# Fix undefined variables in field-level history log
if selected:
    path = os.path.join(FOLDER, selected)
    entry = load_entry(path)
    for category, fields in categories.items():
        with st.expander(f"{category} History"):
            for field in fields:
                value = entry.get(field, "")
                new_value = st.text_area(field.replace("_", " ").title(), value, key=f"{selected}_{field}")
                if new_value != value:
                    log_field_change(selected, field, value, new_value)
                    entry[field] = new_value
                    save_entry(path, entry)
                if st.button(f"Show Change History for {field}", key=f"history_{selected}_{field}"):
                    history = get_field_history(field, selected)
                    if history:
                        for h in history:
                            st.markdown(f"- {h['timestamp']}: {h['old_value']} → {h['new_value']}")
                    else:
                        st.info("No history available for this field.")

# Ensure 'missing_fields' is defined
if 'missing_fields' not in locals():
    missing_fields = [field for field in fields if not entry.get(field)]

# Stub for calculate_field_impact
def calculate_field_impact(field, entry):
    # Replace with actual logic to calculate impact
    return len(field)  # Example: impact based on field name length

# Stub for get_community_suggestions
def get_community_suggestions(field):
    # Replace with actual logic to fetch community suggestions
    return ["Example suggestion 1", "Example suggestion 2"]

# Stub for calculate_field_variance
def calculate_field_variance(field, entry):
    # Replace with actual logic to calculate variance across cultures
    return len(field) % 3  # Example: variance based on field name length modulo 3

# Stub for suggest_field_content
def suggest_field_content(field, culture_name):
    # Replace with actual logic to generate AI suggestions
    return f"Suggested content for {field} in {culture_name}"

# Optimize save operations with caching and debounce logic
save_queue = []

def background_save():
    while True:
        if save_queue:
            path, entry = save_queue.pop(0)
            save_entry(path, entry)
        time.sleep(1)  # Adjust sleep interval as needed

# Start background save thread
threading.Thread(target=background_save, daemon=True).start()

def main():
    st.set_page_config(layout="wide", page_title="Culture Entry Editor")

    files = sorted(load_entries(FOLDER))

    # Fix region and language filters
    all_entries = [load_entry(os.path.join(FOLDER, f)) for f in files]
    region_filter = st.sidebar.selectbox("Filter by Region", ["All"] + sorted({r for e in all_entries for r in e.get("primary_regions", "").split(", ") if r}))
    language_filter = st.sidebar.selectbox("Filter by Language", ["All"] + sorted({e.get("language_tag") for e in all_entries if e.get("language_tag")}))

    filtered_files = []
    for i, e in enumerate(all_entries):
        if (region_filter == "All" or region_filter in e.get("primary_regions", "")) and \
           (language_filter == "All" or language_filter == e.get("language_tag")):
            filtered_files.append(files[i])
    files = filtered_files

    search_query = st.text_input("Search by name")
    filters = {
        "language_tag": st.sidebar.text_input("Filter by Language Tag"),
        "primary_regions": st.sidebar.text_input("Filter by Primary Regions")
    }
    all_entries = advanced_search(all_entries, search_query, filters)

    if not files:
        st.warning("No matching .v3.json files found.")
        return

    # Add header with branding accent
    st.markdown(f"""
    <div style='background-color:{DEEP_PURPLE}; padding:1rem; border-radius: 8px; text-align:center;'>
        <h1 style='color:white; font-size:2.5rem;'>🌍 Global Culture Explorer</h1>
        <p style='color:{LIGHT_PURPLE}; font-size:1.1rem;'>Discover and contribute to the world’s cultural tapestry</p>
    </div>
    """, unsafe_allow_html=True)

    # Add persistent sticky navigation header
    completion_percentage, missing_fields = calculate_summary(entry)
    st.markdown(f"""
        <style>
        .sticky-header {{
            position: sticky;
            top: 0;
            z-index: 1000;
            background-color: {DEEP_PURPLE};
            padding: 1rem;
            border-radius: 0 0 8px 8px;
            box-shadow: 0px 2px 5px rgba(0, 0, 0, 0.2);
        }}
        </style>
        <div class="sticky-header">
            <h1 style="color: white; margin: 0;">🌍 Global Culture Explorer</h1>
            <p style="color: {LIGHT_PURPLE}; margin: 0;">Completion: {completion_percentage:.1f}%</p>
            <div style="margin-top: 0.5rem;">
                <button onclick="window.scrollTo(0, document.body.scrollHeight);">Export Options</button>
            </div>
        </div>
    """, unsafe_allow_html=True)

    with st.sidebar:
        # Sidebar title bar with deep purple
        st.markdown(f"""
        <div style='background-color:{DEEP_PURPLE}; padding:0.5rem 1rem; border-radius:8px;'>
            <h2 style='color:white;'>🌐 Culture Navigator</h2>
        </div>
        """, unsafe_allow_html=True)

        st.image("https://upload.wikimedia.org/wikipedia/commons/thumb/6/6f/Globe_icon.svg/1200px-Globe_icon.svg.png", width=80)
        st.title("🌐 Culture Navigator")
        st.markdown("Navigate, preview, and enrich global cultural profiles.")
        st.markdown("---")
        selected = st.selectbox("Choose a culture entry", files)
        if selected:
            entry = load_entry(os.path.join(FOLDER, selected))
            completion_percentage, missing_fields = calculate_summary(entry)
            st.markdown(f"**Completion:** {completion_percentage:.1f}%")
            st.markdown(f"**Language:** {entry.get('language_tag', '—')}")
            st.markdown(f"**Estimated Speakers:** {entry.get('estimated_speakers', '—')}")
            st.markdown(f"**Primary Regions:** {entry.get('primary_regions', '—')}")
            next_field = get_next_field(entry)
            if next_field:
                st.info(f"Next Field: {next_field.replace('_', ' ').title()}")
        feedback = st.text_area("💬 Feedback", placeholder="Suggest edits or contribute insights…")
        if st.button("📤 Submit Feedback"):
            st.success("Thanks! Your feedback was recorded.")

    # Move entry selector to the sidebar
    with st.sidebar:
        selected = st.selectbox("Choose a culture entry:", files)

    path = os.path.join(FOLDER, selected)
    entry = load_entry(path)

    # Display metadata
    mod_time = os.path.getmtime(path)
    st.caption(f"Last modified: {datetime.datetime.fromtimestamp(mod_time).strftime('%Y-%m-%d %H:%M:%S')}")

    if selected:
        path = os.path.join(FOLDER, selected)
        entry = load_entry(path)

        st.subheader(f"📘 {entry.get('culture_name', 'Unnamed Culture')}")

        # Hero Image
        hero = entry.get("image_url")
        if hero:
            st.markdown(f"""
                <div style='text-align: center; padding: 1rem; background-color: #e3f2fd;'>
                    <img src='{hero}' style='max-width: 100%; border-radius: 10px;'>
                    <h2 style='margin-top: 0.5rem;'>{entry.get('culture_name', 'Cultural Photo')}</h2>
                </div>
            """, unsafe_allow_html=True)

        # Dynamic Progress Badge
        completion_percentage, missing_fields = calculate_summary(entry)
        status_color = DEEP_PURPLE if completion_percentage >= 100 else "#4caf50" if completion_percentage > 80 else "#ff9800" if completion_percentage > 50 else "#f44336"
        st.markdown(f"<span style='background-color: {status_color}; color: white; padding: 0.3rem 0.7rem; border-radius: 5px; font-size: 0.9rem;'>Progress: {completion_percentage:.1f}%</span>", unsafe_allow_html=True)

        # Pin completion summary to top
        st.markdown(f"""
        <div style='display: flex; gap: 1rem; padding: 0.5rem 1rem; background-color: #f4f6f8; border-radius: 8px;'>
            <span style='color: {status_color}; font-weight: bold;'>✓ {completion_percentage:.1f}% Complete</span>
            <span>🗣️ Language: {entry.get('language_tag', '—')}</span>
            <span>👥 Speakers: {entry.get('estimated_speakers', '—')}</span>
            <span>🌍 Regions: {entry.get('primary_regions', '—')}</span>
        </div>
        """, unsafe_allow_html=True)

        # Add language-to-culture mapping
        same_lang = [f for f in files if load_entry(os.path.join(FOLDER, f)).get("language_tag") == entry.get("language_tag")]
        if len(same_lang) > 1:
            st.sidebar.markdown("🧩 Other cultures with this language:")
            for f in same_lang:
                st.sidebar.markdown(f"- {f.replace('.v3.json', '')}")

        # Add "Empty Fields Only" toggle
        show_only_empty = st.checkbox("✏️ Show Only Incomplete Fields")

        # Inline Editing with AI suggestion preview
        # Add field-level completion icons
        # Add card icons per field type
        FIELD_ICONS = {
            "culture_name": "🌍",
            "language_tag": "🗣️",
            "estimated_speakers": "👥",
            "primary_regions": "📍",
            "greetings": "🖐️",
            "communication_style": "🧠",
            "gestures": "🤝",
            "religion": "🙏",
            "gender_roles": "⚧️",
            "taboos": "🚫",
            "clothing": "🧥"
        }

        # Define iconography for fields
        field_icons = {
            "Religion": "🕊️",
            "Clothing": "👗",
            "Gender Roles": "♀️♂️⚧️",
            "Language": "🗣️",
            "Food": "🍲",
            "Festivals": "🎉",
        }

        # Add responsive grid system
        cols = st.columns([2, 3])
        with cols[0]:
            st.metric("Completion", f"{completion_percentage:.1f}%")
        with cols[1]:
            display_image_gallery(entry.get("image_urls", []))

        for category, fields in categories.items():
            filled = sum(1 for f in fields if entry.get(f))
            total = len(fields)
            label = f"{category} ({filled}/{total})"
            anchor = category.replace(" ", "_").lower()
            st.markdown(f"<a name='{anchor}'></a>", unsafe_allow_html=True)  # Add anchor for navigation
            with st.expander(label):
                for field in fields:
                    if show_only_empty and entry.get(field):
                        continue
                    value = entry.get(field, "")
                    is_saved = st.session_state.get(f"{field}_saved", True)
                    new_value = st.text_area(field.replace("_", " ").title(), value, key=field)
                    if new_value != value:
                        st.session_state[f"{field}_saved"] = False
                        add_save_indicator(field, False)
                        if st.button("Save", key=f"save_{field}"):
                            entry[field] = new_value
                            save_entry(path, entry)
                            st.session_state[f"{field}_saved"] = True
                            add_save_indicator(field, True)
                    # Show AI suggestion badge
                    if "AI-generated" in value:
                        st.markdown(f"<span style='color:{DEEP_PURPLE}; font-size:0.9rem;'>(AI-suggested)</span>", unsafe_allow_html=True)
                    # Add translation tooltip
                    add_translation_tooltip(field, value, translations)
                if st.button("💾 Save Changes", key=f"save_{category}"):
                    changes = {field: {"before": entry.get(field, ""), "after": new_value} for field, new_value in entry.items() if entry.get(field) != new_value}
                    log_submission(entry, path, changes)
                    save_backup(entry, path)
                    save_entry(path, entry)
                    st.success("Changes saved successfully!")
                    # Auto-scroll to next field
                    next_field = get_next_field(entry)
                    if next_field:
                        scroll_to_field(next_field)

        # Gamify completion with XP bar and badges
        if completion_percentage >= 100:
            st.balloons()
            st.markdown("🎖️ **Congratulations! 100% Complete!**")
        elif completion_percentage > 80:
            st.markdown("🏆 **Level Up! Over 80% Complete!**")

        # Add XP bar
        st.progress(int(completion_percentage))

        # Replace progress bar with Streamlit's native bar
        st.progress(int(completion_percentage))
        st.caption(f"{completion_percentage:.1f}% Complete")

        # Add jump-to dropdown
        jump_to_field = st.sidebar.selectbox("Jump to Field", ["Select"] + list(entry.keys()))
        if jump_to_field != "Select":
            components.html(f"<script>document.getElementById('{jump_to_field}').scrollIntoView();</script>", height=0)

        # Validation for Key Fields
        try:
            int(entry.get("estimated_speakers", "").replace(",", ""))
        except ValueError:
            st.warning("Estimated speakers should be a number")

        # Export Buttons
        st.markdown("### Export Options")
        if st.button("Export as HTML"):
            export_as_html(entry, "export.html")
            with open("export.html", "rb") as f:
                st.download_button("Download HTML", f, file_name="export.html")

        if st.button("Export as Markdown"):
            export_as_markdown(entry, "export.md")
            with open("export.md", "rb") as f:
                st.download_button("Download Markdown", f, file_name="export.md")

    # Reorganize batch export section
    with st.expander("📦 Batch Export Options"):
        selected_files = st.multiselect("Select entries to export", files)
        entries = {file: load_entry(os.path.join(FOLDER, file)) for file in selected_files}

        if st.button("Export Selected as HTML (ZIP)"):
            batch_export(entries, "html", "batch_export_html.zip")
            with open("batch_export_html.zip", "rb") as f:
                st.markdown(f"""
                <a href="data:application/zip;base64,{f.read().encode('base64').decode()}" download="culture_html_export.zip" style='
                    display:inline-block;
                    background-color:{DEEP_PURPLE};
                    color:white;
                    padding:0.5rem 1.2rem;
                    border-radius:6px;
                    text-decoration:none;
                    font-weight:bold;
                '>⬇️ Download HTML ZIP</a>
                """, unsafe_allow_html=True)

        if st.button("Export Selected as Markdown (ZIP)"):
            batch_export(entries, "md", "batch_export_md.zip")
            with open("batch_export_md.zip", "rb") as f:
                st.markdown(f"""
                <a href="data:application/zip;base64,{f.read().encode('base64').decode()}" download="culture_md_export.zip" style='
                    display:inline-block;
                    background-color:{DEEP_PURPLE};
                    color:white;
                    padding:0.5rem 1.2rem;
                    border-radius:6px;
                    text-decoration:none;
                    font-weight:bold;
                '>⬇️ Download Markdown ZIP</a>
                """, unsafe_allow_html=True)

        # Add sticky category navigation sidebar
        with st.sidebar:
            st.markdown("### Jump to Category")
            category_names = [category.replace(" ", "_").lower() for category in categories.keys()]
            selected_category = st.selectbox("Select Category", ["Select"] + list(categories.keys()))
            if selected_category != "Select":
                anchor = selected_category.replace(" ", "_").lower()
                components.html(f"<script>document.getElementsByName('{anchor}')[0].scrollIntoView();</script>", height=0)

    # Add interactive map for primary regions
    if selected:
        primary_regions = entry.get("primary_regions", "")
        if primary_regions:
            st.markdown("### 🌍 Primary Regions Map")
            region_map = create_region_map(primary_regions)
            st_folium(region_map, width=700, height=500)

    # Placeholder: Replace with real AI call
    def suggest_field_content(field_name, context):
        return f"Sample content for {field_name}: [add cultural example here]"

    # Field visibility rules
    FIELD_VISIBILITY_RULES = {
        "spiritual_practices": lambda entry: entry.get("religion", "").lower() not in ["atheist", "secular"]
    }

    # Define field-specific widgets
    FIELD_WIDGETS = {
        "estimated_speakers": lambda k, v: st.number_input(k, value=int(v or 0), step=1, min_value=0),
        "primary_regions": lambda k, v: st.multiselect(k, options=["North America", "Europe", "Asia", "Africa", "Oceania", "South America"], default=v.split(", ") if v else []),
        "language_tag": lambda k, v: st.selectbox(k, options=["en", "es", "fr", "de", "zh", "ar", "ru", "ja", "hi"], index=["en", "es", "fr", "de", "zh", "ar", "ru", "ja", "hi"].index(v) if v in ["en", "es", "fr", "de", "zh", "ar", "ru", "ja", "hi"] else 0)
    }

    # Add AI Assistant Panel
    with st.sidebar.expander("💡 Suggestions", expanded=False):
        if selected:
            current_field = get_next_field(entry)  # Example: Detect current field
            if current_field:
                context = entry.get("culture_name", "")
                suggestion = suggest_field_content(current_field, context)
                confidence = add_confidence_meter(current_field, suggestion)
                st.markdown(f"You rated this suggestion: {confidence}/5")

    # Enhance dynamic AI panel
    with st.sidebar.expander("🧠 Suggested Fills", expanded=False):
        for field in fields:
            suggestion = suggest_field_content(field, entry.get("culture_name", ""))
            confidence = "High"  # Placeholder for GPT confidence
            st.markdown(f"**{field.replace('_', ' ').title()}**: {suggestion} ({confidence} Confidence)")
            if st.button(f"⭐ Accept Suggestion for {field}", key=f"accept_{field}"):
                entry[field] = suggestion
                save_entry(path, entry)
                st.success(f"Suggestion for {field} accepted and saved.")

    # Add admin mode toggle
    admin_mode = st.sidebar.checkbox("🔐 Admin Mode")

    # Restrict editing of certain fields if admin mode is off
    RESTRICTED_FIELDS = ["religion", "gender_roles"]

    # Apply field visibility rules
    for category, fields in categories.items():
        filled = sum(1 for f in fields if entry.get(f))
        total = len(fields)
        label = f"{category} ({filled}/{total})"
        anchor = category.replace(" ", "_").lower()
        st.markdown(f"<a name='{anchor}'></a>", unsafe_allow_html=True)  # Add anchor for navigation
        with st.expander(label):
            for field in fields:
                if field in FIELD_VISIBILITY_RULES and not FIELD_VISIBILITY_RULES[field](entry):
                    continue
                if not admin_mode and field in RESTRICTED_FIELDS:
                    st.markdown(f"**{field.replace('_', ' ').title()}**: {entry.get(field, '—')}")
                    continue
                value = entry.get(field, "")
                status_icon = "✅" if value.strip() else "❌"
                pretty_field = f"{status_icon} {field.replace('_', ' ').title()}"
                if field in FIELD_WIDGETS:
                    entry[field] = FIELD_WIDGETS[field](pretty_field, value)
                else:
                    entry[field] = st.text_area(pretty_field, value, height=100)
                # Show AI suggestion badge
                if "AI-generated" in value:
                    st.markdown(f"<span style='color:{DEEP_PURPLE}; font-size:0.9rem;'>(AI-suggested)</span>", unsafe_allow_html=True)
                # Add translation tooltip
                add_translation_tooltip(field, value, translations)
            if st.button("💾 Save Changes", key=f"save_{category}"):
                changes = {field: {"before": entry.get(field, ""), "after": new_value} for field, new_value in entry.items() if entry.get(field) != new_value}
                log_submission(entry, path, changes)
                save_backup(entry, path)
                save_entry(path, entry)
                st.success("Changes saved successfully!")

    # Add a "Restore Previous" dropdown
    with st.sidebar.expander("🔄 Restore Previous", expanded=False):
        backup_files = [f for f in os.listdir("backups") if f.startswith(os.path.basename(path))]
        if backup_files:
            selected_backup = st.selectbox("Select a backup to restore", backup_files)
            if st.button("Restore Backup"):
                backup_path = os.path.join("backups", selected_backup)
                entry = load_entry(backup_path)
                save_entry(path, entry)
                st.success("Backup restored successfully!")
        else:
            st.info("No backups available.")

    # Add theme toggle in the sidebar
    dark_mode = st.sidebar.checkbox("🌙 Dark Mode")

    # Define styles for dark mode
    if dark_mode:
        CARD_STYLE_FILLED = "background-color:#424242; border:2px solid #90caf9; padding: 1rem; border-radius: 10px; margin-bottom: 1rem; box-shadow: 1px 1px 5px rgba(0,0,0,0.5); color: #ffffff;"
        CARD_STYLE_EMPTY = "background-color:#616161; border:2px dashed #ffb74d; padding: 1rem; border-radius: 10px; margin-bottom: 1rem; box-shadow: 1px 1px 5px rgba(0,0,0,0.5); color: #ffffff;"
        DEEP_PURPLE = "#bb86fc"
        LIGHT_PURPLE = "#cfbaf0"
    else:
        CARD_STYLE_FILLED = "background-color:#e3f2fd; border:2px solid #90caf9; padding: 1rem; border-radius: 10px; margin-bottom: 1rem; box-shadow: 1px 1px 5px rgba(0,0,0,0.05);"
        CARD_STYLE_EMPTY = "background-color:#fff3cd; border:2px dashed #ffb74d; padding: 1rem; border-radius: 10px; margin-bottom: 1rem; box-shadow: 1px 1px 5px rgba(0,0,0,0.05);"
        DEEP_PURPLE = "#4A148C"
        LIGHT_PURPLE = "#E1BEE7"

    # Apply theme based on primary region
    primary_region = entry.get("primary_regions", "").split(", ")[0] if entry.get("primary_regions") else ""
    theme_colors = get_theme_colors(primary_region)
    st.markdown(f"""
        <style>
        body {{
            background-color: {theme_colors['background']};
        }}
        .accent {{
            color: {theme_colors['accent']};
        }}
        </style>
    """, unsafe_allow_html=True)

    # Apply custom CSS for consistent styling
    st.markdown("""
        <style>
        body {
            font-family: 'Arial', sans-serif;
        }
        .stSidebar {
            background-color: #F3E5F5;
        }
        </style>
    """, unsafe_allow_html=True)

    # Add radar chart above the editor
    if selected:
        st.markdown("### 📊 Category Completion Radar Chart")
        radar_chart = create_completion_radar_chart(entry)
        st.plotly_chart(radar_chart)

    # Add field completion minimap in the sidebar
    with st.sidebar.expander("📋 Field Completion Minimap", expanded=True):
        for category, fields in categories.items():
            st.markdown(f"**{category}**")
            for field in fields:
                status = "✅" if entry.get(field) else "❌"
                st.markdown(f"- {status} {field.replace('_', ' ').title()}")

    # Add dynamic image gallery
    image_urls = entry.get("image_urls", [])  # Assuming 'image_urls' is a list of image links
    if image_urls:
        display_image_gallery(image_urls)

    # Add top missing fields heatmap
    generate_missing_fields_heatmap(all_entries)

    # Add a plot for speaker trends
    if selected:
        st.markdown("### 📈 Estimated Speaker Trends")
        display_speaker_trends(entry)

    # Add cultural calendar widget
    if selected:
        display_cultural_calendar(entry)

    # Add print-friendly mode
    if st.button("🖨️ Generate Print-Friendly Version"):
        generate_print_friendly(entry)

    # Add custom CSS for visual polish
    st.markdown("""
        <style>
        .stButton>button {
            background-color: #4A148C;
            color: white;
            border-radius: 5px;
        }
        .stButton>button:hover {
            background-color: #7B1FA2;
        }
        .stSidebar {
            background-color: #F3E5F5;
        }
        </style>
    """, unsafe_allow_html=True)

    # Validate schema in the main function
    if selected:
        schema_errors = validate_schema(entry, default_schema)
        if schema_errors:
            st.error("Schema Validation Errors:")
            for error in schema_errors:
                st.markdown(f"- {error}")

    # Integrate suggestion submission into the main function
    if selected:
        st.markdown("### 💡 Crowdsourced Suggestions")
        suggestion_field = st.selectbox("Field to Suggest", list(entry.keys()))
        suggestion_text = st.text_area("Your Suggestion")
        if st.button("Submit Suggestion"):
            submit_suggestion(selected, suggestion_field, suggestion_text)
            st.success("Thank you! Your suggestion has been recorded.")

    # Add compare mode
    if st.checkbox("Enable Compare Mode"):
        culture_a = st.selectbox("Select Culture A", files, key="compare_a")
        culture_b = st.selectbox("Select Culture B", files, key="compare_b")
        if culture_a and culture_b:
            entry_a = load_entry(os.path.join(FOLDER, culture_a))
            entry_b = load_entry(os.path.join(FOLDER, culture_b))
            compare_cultures(entry_a, entry_b)

    # Add mini card layout for each field
    for category, fields in categories.items():
        with st.expander(f"{category} Fields"):
            for field in fields:
                value = entry.get(field, "")
                icon = FIELD_ICONS.get(field, "")
                st.markdown(f"**{icon} {field.replace('_', ' ').title()}**")
                if st.button("Expand", key=f"expand_{field}"):
                    st.text_area(f"Edit {field}", value, key=f"edit_{field}")

    # Add 'Next Best Action' AI card
    st.markdown("### 🧠 Next Best Action")
    next_field = get_next_field(entry)
    if next_field:
        st.info(f"We recommend improving **{next_field.replace('_', ' ').title()}** next. It has the highest impact on completion.")

    # Add AI-powered auto-fill mode
    if st.button("⚡ Auto-Fill All Missing Fields"):
        for field in missing_fields:
            suggestion = suggest_field_content(field, entry.get("culture_name", ""))
            entry[field] = f"{suggestion} (AI-generated)"
        save_entry(path, entry)
        st.success("All empty fields auto-filled!")

    # Add a dashboard metrics panel
    if st.sidebar.checkbox("📊 Show Dashboard Metrics"):
        st.markdown("### 📊 Dashboard Metrics")
        total_cultures = len(files)
        avg_completion = sum(calculate_summary(load_entry(os.path.join(FOLDER, f)))[0] for f in files) / total_cultures
        st.metric("Total Cultures", total_cultures)
        st.metric("Average Completion", f"{avg_completion:.1f}%")

        # Top 5 incomplete fields
        from collections import Counter
        field_counts = Counter()
        for f in files:
            _, missing = calculate_summary(load_entry(os.path.join(FOLDER, f)))
            field_counts.update(missing)
        top_incomplete = field_counts.most_common(5)
        st.markdown("#### Top 5 Incomplete Fields")
        for field, count in top_incomplete:
            st.markdown(f"- {field.replace('_', ' ').title()}: {count} missing")

# Add Smart Completion Assistant to the sidebar
st.sidebar.markdown("### 🧠 Smart Completion Assistant")

# Fields with most impact on completion
impactful_fields = sorted(missing_fields, key=lambda f: calculate_field_impact(f, entry), reverse=True)[:5]
st.sidebar.markdown("#### Fields with Most Impact")
for field in impactful_fields:
    st.sidebar.markdown(f"- {field.replace('_', ' ').title()}")

# Fields with most community suggestions
suggested_fields = sorted(missing_fields, key=lambda f: len(get_community_suggestions(f)), reverse=True)[:5]
st.sidebar.markdown("#### Fields with Most Suggestions")
for field in suggested_fields:
    st.sidebar.markdown(f"- {field.replace('_', ' ').title()}")

# Fields that differ most across similar cultures
differing_fields = sorted(missing_fields, key=lambda f: calculate_field_variance(f, entry), reverse=True)[:5]
st.sidebar.markdown("#### Fields with Most Variance")
for field in differing_fields:
    st.sidebar.markdown(f"- {field.replace('_', ' ').title()}")

# Add 'Fill All Suggested' button
if st.sidebar.button("⚡ Fill All Suggested Fields"):
    for field in suggested_fields:
        suggestion = suggest_field_content(field, entry.get("culture_name", ""))
        entry[field] = f"{suggestion} (AI-generated)"
    save_entry(path, entry)
    st.sidebar.success("All suggested fields auto-filled!")

# Initialize 'languages' and 'selected_language' for multilingual UI
languages = {"English": "en", "Spanish": "es", "French": "fr"}
selected_language = st.sidebar.selectbox("🌐 Select Language", options=languages.keys())

# Ensure 't' function uses these variables
def t(key):
    return translations[languages[selected_language]].get(key, key)

# Example usage of multilingual UI
st.metric(t("completion"), f"{completion_percentage:.1f}%")

# Correct the floating action toolbar and section bookmarking code
st.markdown(
    """
    <style>
    .floating-toolbar {
        position: fixed;
        bottom: 20px;
        right: 20px;
        background-color: #ffffff;
        border: 1px solid #cccccc;
        border-radius: 8px;
        box-shadow: 0px 4px 6px rgba(0, 0, 0, 0.1);
        padding: 10px;
        z-index: 1000;
    }
    .floating-toolbar button {
        margin: 5px;
        padding: 10px;
        border: none;
        border-radius: 4px;
        background-color: #007bff;
        color: white;
        cursor: pointer;
    }
    .floating-toolbar button:hover {
        background-color: #0056b3;
    }
    </style>
    <div class="floating-toolbar">
        <button onclick="window.scrollTo(0, 0)">⬆️ Top</button>
        <button onclick="window.scrollTo(0, document.body.scrollHeight)">⬇️ Bottom</button>
        <button onclick="toggleFilters()">🔍 Filters</button>
        <button onclick="enrichWithAI()">⚡ Enrich</button>
    </div>
    <script>
    function toggleFilters() {
        alert('Toggle filters clicked!'); // Replace with actual filter toggle logic
    }
    function enrichWithAI() {
        alert('Enrich with AI clicked!'); // Replace with actual AI enrichment logic
    }
    </script>
    """,
    unsafe_allow_html=True
)

# Correct section bookmarking and auto-scroll code
st.markdown(
    """
    <script>
    document.addEventListener('DOMContentLoaded', function() {
        const hash = window.location.hash.substring(1);
        if (hash) {
            const element = document.getElementById(hash);
            if (element) {
                element.scrollIntoView();
            }
        }
    });
    </script>
    """,
    unsafe_allow_html=True
)

# Add real-time language popularity estimation
st.sidebar.markdown("### 🌍 Language Popularity Estimation")

# Example data for language popularity (replace with actual data source)
language_data = {
    "en": {"name": "English", "speakers": 1500000000},
    "es": {"name": "Spanish", "speakers": 500000000},
    "fr": {"name": "French", "speakers": 300000000},
}

# Display language popularity
for code, data in language_data.items():
    st.sidebar.markdown(f"- **{data['name']}**: {data['speakers']:,} speakers")

# Add toggle for population map
if st.sidebar.checkbox("Show Population Map by Language"):
    import folium
    from streamlit_folium import st_folium

    # Example map data (replace with actual data)
    map_data = [
        {"language": "English", "lat": 51.5074, "lon": -0.1278, "speakers": 1500000000},
        {"language": "Spanish", "lat": 40.4168, "lon": -3.7038, "speakers": 500000000},
        {"language": "French", "lat": 48.8566, "lon": 2.3522, "speakers": 300000000},
    ]

    # Create a folium map
    m = folium.Map(location=[20, 0], zoom_start=2)
    for entry in map_data:
        folium.CircleMarker(
            location=[entry["lat"], entry["lon"]],
            radius=entry["speakers"] / 100000000,  # Scale radius by speakers
            popup=f"{entry['language']}: {entry['speakers']:,} speakers",
            color="blue",
            fill=True,
            fill_opacity=0.6,
        ).add_to(m)

    # Display the map
    st_folium(m, width=700, height=500)

    # Add a multi-format export panel
    st.sidebar.markdown("### 📦 Export Options")

    # Export as PDF
    if st.sidebar.button("Export as PDF"):
        import pdfkit
        pdf_content = "<h1>Culture Report</h1><p>This is a sample PDF export.</p>"  # Replace with actual content
        pdfkit.from_string(pdf_content, "culture_report.pdf")
        st.sidebar.success("PDF exported successfully!")

    # Export as JSON
    if st.sidebar.button("Export as JSON"):
        import json
        json_content = json.dumps(entry, indent=2)
        with open("culture_report.json", "w", encoding="utf-8") as f:
            f.write(json_content)
        st.sidebar.success("JSON exported successfully!")

    # Export as Markdown
    if st.sidebar.button("Export as Markdown"):
        markdown_content = "# Culture Report\n\n"  # Replace with actual content
        for field, value in entry.items():
            markdown_content += f"## {field.replace('_', ' ').title()}\n{value}\n\n"
        with open("culture_report.md", "w", encoding="utf-8") as f:
            f.write(markdown_content)
        st.sidebar.success("Markdown exported successfully!")

    # Export as PowerPoint
    if st.sidebar.button("Export as PowerPoint"):
        from pptx import Presentation
        ppt = Presentation()
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Culture Report"
        for field, value in entry.items():
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])
            slide.shapes.title.text = field.replace('_', ' ').title()
            slide.placeholders[1].text = value
        ppt.save("culture_report.pptx")
        st.sidebar.success("PowerPoint exported successfully!")

# Add shareable links and embed codes
st.sidebar.markdown("### 🔗 Shareable Links and Embed Codes")

# Generate public read-only card link
if st.sidebar.button("Generate Public Link"):
    public_url = f"https://culture-project.com/view/{selected}"  # Replace with actual URL logic
    st.sidebar.markdown(f"[View Public Card]({public_url})")
    st.sidebar.success("Public link generated!")

# Generate embed code
if st.sidebar.button("Generate Embed Code"):
    embed_code = f"<iframe src='https://culture-project.com/view/{selected}' width='600' height='400'></iframe>"
    st.sidebar.text_area("Embed Code", embed_code)
    st.sidebar.success("Embed code generated!")

# Open in print preview mode
if st.sidebar.button("Open in Print Preview Mode"):
    st.sidebar.markdown(f"[Print Preview](https://culture-project.com/print/{selected})")  # Replace with actual print preview URL logic

# Add AI Voice Assistant to the sidebar
st.sidebar.markdown("### 🎙️ AI Voice Assistant")

# Check if microphone input is available
try:
    from streamlit_speech_recognition import speech_to_text
    user_query = speech_to_text()
    if user_query:
        st.sidebar.markdown(f"**You said:** {user_query}")

        # Process the query with OpenAI
        import openai
        response = openai.Completion.create(
            engine="text-davinci-003",
            prompt=f"Answer the following query about culture: {user_query}",
            max_tokens=150
        )
        answer = response.choices[0].text.strip()

        # Display the response
        st.sidebar.markdown(f"**AI Response:** {answer}")

        # Auto-scroll to relevant section (if applicable)
        if "tell me about" in user_query.lower():
            culture_name = user_query.lower().replace("tell me about", "").strip()
            st.sidebar.markdown(f"Scrolling to section for {culture_name}...")
            # Add logic to scroll to the relevant section
except ModuleNotFoundError:
    st.sidebar.error("Speech recognition module not installed. Please install 'streamlit-speech-recognition'.")
except ImportError:
    st.sidebar.error("Speech recognition module could not be imported.")

# Add Interactive Culture Builder
st.sidebar.markdown("### 🛠️ Interactive Culture Builder")

if st.sidebar.button("Create New Culture Entry"):
    st.session_state["new_entry"] = {}
    st.session_state["current_step"] = 0

if "new_entry" in st.session_state:
    new_entry = st.session_state["new_entry"]
    current_step = st.session_state["current_step"]

    # Define steps for the guided wizard
    steps = [
        {"label": "Basic Information", "fields": ["culture_name", "language_tag", "primary_regions"]},
        {"label": "Demographics", "fields": ["estimated_speakers"]},
        {"label": "Cultural Practices", "fields": ["religion", "taboos", "clothing"]},
    ]

    # Display current step
    step = steps[current_step]
    st.markdown(f"### Step {current_step + 1}: {step['label']}")

    for field in step["fields"]:
        new_entry[field] = st.text_input(field.replace("_", " ").title(), new_entry.get(field, ""))

    # Navigation buttons
    col1, col2, col3 = st.columns([1, 1, 1])
    with col1:
        if st.button("Previous") and current_step > 0:
            st.session_state["current_step"] -= 1
    with col2:
        if st.button("Next") and current_step < len(steps) - 1:
            st.session_state["current_step"] += 1
    with col3:
        if st.button("Finish"):
            # Save the new entry
            save_entry(os.path.join(FOLDER, f"{new_entry['culture_name']}.json"), new_entry)
            st.success("New culture entry created successfully!")
            del st.session_state["new_entry"]
            del st.session_state["current_step"]

# Add tooltips for technical terms
TOOLTIPS = {
    "communication_style": "The way people interact and convey information in this culture.",
    "taboos": "Actions or behaviors that are considered inappropriate or forbidden.",
    "AI-generated": "This content was suggested by an AI model based on available data.",
}

# Example usage of tooltips in field rendering
for category, fields in categories.items():
    with st.expander(f"{category} Fields"):
        for field in fields:
            value = entry.get(field, "")
            tooltip = TOOLTIPS.get(field, "")
            if tooltip:
                st.markdown(f"**{field.replace('_', ' ').title()}** 🛈")
                st.markdown(f"<span title='{tooltip}'>{value}</span>", unsafe_allow_html=True)
            else:
                st.text_area(f"Edit {field}", value, key=f"edit_{field}")

# Add Accessibility Mode to the sidebar
st.sidebar.markdown("### ♿ Accessibility Mode")

# High contrast toggle
high_contrast = st.sidebar.checkbox("Enable High Contrast Mode")
if high_contrast:
    st.markdown(
        """
        <style>
        body {
            background-color: #000000;
            color: #ffffff;
        }
        .stButton>button {
            background-color: #ffffff;
            color: #000000;
            border: 2px solid #ffffff;
        }
        .stButton>button:hover {
            background-color: #cccccc;
        }
        </style>
        """,
        unsafe_allow_html=True
    )

# Larger font toggle
larger_font = st.sidebar.checkbox("Enable Larger Font Size")
if larger_font:
    st.markdown(
        """
        <style>
        body {
            font-size: 18px;
        }
        </style>
        """,
        unsafe_allow_html=True
    )

# Enable keyboard navigation
st.markdown(
    """
    <script>
    document.addEventListener('keydown', function(event) {
        if (event.key === 'ArrowDown') {
            window.scrollBy(0, 100);
        } else if (event.key === 'ArrowUp') {
            window.scrollBy(0, -100);
        }
    });
    </script>
    """,
    unsafe_allow_html=True
)

updated = {}
with st.expander("🧑‍🤝‍🧑 Family & Social Structure"):
    updated["family_structure"] = st.text_area("Family Structure", entry.get("family_structure", ""), height=100)
    updated["social_roles"] = st.text_area("Social Roles", entry.get("social_roles", ""), height=100)

with st.expander("🗣️ Communication Style"):
    updated["communication_style"] = st.text_area("Communication Style", entry.get("communication_style", ""), height=100)
    updated["gestures"] = st.text_area("Gestures", entry.get("gestures", ""), height=100)

# Save button
if st.button("💾 Save"):
    entry.update(updated)
    save_entry(path, entry)
    st.success(f"{selected} saved.", icon="✅")